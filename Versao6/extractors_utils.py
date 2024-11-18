from docx import Document
import fitz  # PyMuPDF
import pytesseract
from PIL import Image, ImageEnhance, ImageOps
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import json
import requests
from io import BytesIO, StringIO
import os
import zipfile
from pptx import Presentation
import pandas as pd
import pdfplumber
import warnings
import tempfile
import rarfile
import py7zr
import sys
import pikepdf
from PIL import ImageFilter
######### Conversão de arquivos .doc para .docx ########
# from docx.document import Document as Document_docx
from win32com.client import Dispatch
# ######## Fim da Conversão de arquivos .doc para .docx ########


sys.path.insert(0, './Versao6')

# Suprime os avisos do tipo UserWarning, incluindo o aviso do openpyxl
warnings.simplefilter("ignore", UserWarning)

# Caminho para o executável do Tesseract no Windows (path)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Caminhos do LibreOffice e Tesseract
libreoffice_path = r"C:\Program Files\LibreOffice\program"

os.environ["PATH"] += os.pathsep + libreoffice_path

tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_xml(file_path_or_content):
    try:
        # Verifica se é conteúdo binário ou um caminho de arquivo
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string
            file_content = file_path_or_content.decode('utf-8')
        else:
            # Lê o arquivo diretamente do caminho fornecido
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = file.read()
        
        # Analisa o conteúdo XML
        root = ET.fromstring(file_content)

        # Função recursiva para extrair texto de todos os elementos XML
        def extract_text(element):
            text = element.text.strip() if element.text else ""
            for child in element:
                text += "\n" + extract_text(child)
            return text
        
        # Extrai o texto do XML
        extracted_text = extract_text(root)
        
        return extracted_text, True, ""  # Retorna o texto extraído e sucesso

    except ET.ParseError as parse_error:
        erro_msg = f"Erro ao analisar o XML: {parse_error}"
        print(erro_msg)
        return "", False, erro_msg
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo XML: {e}"
        print(erro_msg)
        return "", False, erro_msg


# Extração de texto de arquivos ODT (LibreOffice Writer)
def extract_text_from_odt(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)

        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        text_elements = [elem.text for elem in root.iter() if elem.text is not None]
        extracted_text = "\n".join(text_elements)
        return extracted_text, True
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo ODT: {e}"
        print(erro_msg)  # Exibe o erro no console
        return "", False


# Extração de texto de arquivos ODS (LibreOffice Calc)
def extract_text_from_ods(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)
        
        # Abre o arquivo 'content.xml' dentro do ODS
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Inicializa uma lista para armazenar o conteúdo extraído das células
        cell_texts = []

        # Percorre os elementos de célula no arquivo ODS
        for cell in root.iter('{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-cell'):
            # Verifica o conteúdo do texto dentro da célula
            text_content = ''.join(cell.itertext())
            if text_content:  # Adiciona o conteúdo apenas se não for vazio
                cell_texts.append(text_content)

        # Junta o conteúdo de todas as células em uma única string
        extracted_text = "\n".join(cell_texts)
        return extracted_text, True
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo ODS: {e}")
        return "", False


# Extração de texto de arquivos ODP (LibreOffice Impress)
def extract_text_from_odp(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)

        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        text_elements = [elem.text for elem in root.iter() if elem.text is not None]
        extracted_text = "\n".join(text_elements)
        return extracted_text, True
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo ODP: {e}"
        print(erro_msg)  # Exibe o erro no console
        return "", False


# Extração de texto de arquivos ODG (LibreOffice Draw)
def extract_text_from_odg(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)
        
        # Abre o arquivo 'content.xml' dentro do ODG
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Inicializa uma lista para armazenar o conteúdo extraído
        extracted_texts = []

        # Itera sobre todos os elementos de texto relevantes no arquivo ODG, incluindo texto dentro de tags aninhadas
        for elem in root.iter():
            # Verifica se o elemento possui texto (ou parte de texto) e o adiciona à lista
            if elem.text:
                extracted_texts.append(elem.text)
            if elem.tail:  # Também captura texto após a tag de fechamento
                extracted_texts.append(elem.tail)

        # Junta o conteúdo extraído em uma única string
        extracted_text = "\n".join(extracted_texts)
        return extracted_text, True
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo ODG: {e}")
        return "", False


def extract_text_from_txt(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            resultado = file_path_or_content.decode('latin1')
        else:
            with open(file_path_or_content, 'r', encoding='latin1') as file:
                resultado = file.read()
        return resultado, True
    except Exception as e:
        print(f"Erro ao ler arquivo txt: {e}")
        return "", False


def extract_text_from_odf(file_path, extension):
    try:
        if extension == 'odt':
            return extract_text_from_odt(file_path), True, None
        elif extension == 'ods':
            return extract_text_from_ods(file_path), True, None
        elif extension == 'odp':
            return extract_text_from_odp(file_path), True, None
        elif extension == 'odg':
            return extract_text_from_odg(file_path), True, None
    except Exception as e:
        erro_msg = f"Erro ao processar arquivo ODF ({extension}): {e}"
        return "", False, erro_msg


def extract_text_from_xlsx(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            file_content = BytesIO(file_path_or_content)
            df = pd.read_excel(file_content, engine='openpyxl')
        else:
            df = pd.read_excel(file_path_or_content, engine='openpyxl')
        return df.to_string(), True
    except Exception as e:
        print(f"Erro ao ler arquivo xlsx: {e}")
        return "", False


def extract_text_from_csv(file_path_or_content):
    try:
        # Verifica se o conteúdo é binário (baixado de uma URL e não de um arquivo local)
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string usando StringIO
            file_content = StringIO(file_path_or_content.decode('utf-8'))
        else:
            # Se for um caminho local, abre o arquivo e lê o conteúdo
            file_content = file_path_or_content

        # Usa pandas para ler o CSV e convertê-lo em uma string
        df = pd.read_csv(file_content, on_bad_lines='skip')
        return df.to_string(), True, ""  # Retorna o conteúdo como string, status de sucesso e erro vazio
    except Exception as e:
        erro_msg = f"Erro ao ler arquivo CSV: {e}"
        print(erro_msg)
        return "", False, erro_msg


def extract_text_and_images_from_docx(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            doc = Document(BytesIO(file_path_or_content))
        else:
            doc = Document(file_path_or_content)

        all_text = ""
        extracted_text_from_images = []
        extracted_tables = []

        for para in doc.paragraphs:
            all_text += para.text + "\n"

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            extracted_tables.append(table_data)

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                text_from_image = extract_text_from_image(image_data)
                if isinstance(text_from_image, str):
                    extracted_text_from_images.append(text_from_image)

        combined_text = all_text + "\n".join(extracted_text_from_images)
        combined_text += "\n\n" + "\n".join(str(table) for table in extracted_tables)

        return combined_text, True

    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo DOCX: {e}")
        return "", False


def download_and_convert_doc_to_docx(file_path):
    output_path = tempfile.mkdtemp()
    try:
        word = Dispatch("Word.Application")
        doc = word.Documents.Open(file_path)
        doc.SaveAs(os.path.join(output_path, "arquivo_temp.docx"), FileFormat=16)  # 16 = wdFormatDocumentDefault
        doc.Close()
        word.Quit()
        word = None
        return os.path.join(output_path, "arquivo_temp.docx")
    except Exception as e:
        print(f"Erro ao converter o arquivo {file_path}: {e}")
        return None
    

def extract_text_from_image(image_path):
    try:
        # Abre a imagem a partir do caminho
        image = Image.open(image_path)

        # Pré-processamento da imagem
        image = image.convert('L').resize((image.width * 2, image.height * 2), resample=Image.BICUBIC)
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(1.5)
        image = image.point(lambda x: 0 if x < 140 else 255, '1')

        # Realizar OCR na imagem
        text = pytesseract.image_to_string(image, lang='por', config='--psm 5')

        return text, True

    except Exception as e:
        print(f"Erro ao realizar OCR na imagem: {e}")
        return "", False


def extract_text_from_html(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            # Se o conteúdo for bytes (Ex: no caso de uma URL), converte para string
            file_content = file_path_or_content.decode('utf-8')
        else:
            # Se for um caminho de arquivo local
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = file.read()

        soup = BeautifulSoup(file_content, 'html.parser')
        # Extrai todo o texto visível do HTML
        text = soup.get_text(separator="\n")
        return text, True, ""  # Adicionei uma string vazia como mensagem de erro
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo HTML: {e}")
        return "", False, str(e)


def extract_text_from_json(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            file_content = json.loads(file_path_or_content.decode('utf-8'))
        else:
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = json.load(file)

        # Recorre ao conteúdo do JSON e converte tudo em uma string
        def extract_json_text(data, result=""):
            if isinstance(data, dict):
                for key, value in data.items():
                    result += f"{key}: {extract_json_text(value)}\n"
            elif isinstance(data, list):
                for item in data:
                    result += extract_json_text(item) + "\n"
            else:
                result += str(data)
            return result

        extracted_text = extract_json_text(file_content)
        return extracted_text, True, ""  # Adicionei uma string vazia como mensagem de erro
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo JSON: {e}")
        return "", False, str(e)


def extract_text_from_pptx(file_path_or_content):

    try:
        if isinstance(file_path_or_content, (bytes, BytesIO)):
            file_content = BytesIO(file_path_or_content) if isinstance(file_path_or_content, bytes) else file_path_or_content
        elif isinstance(file_path_or_content, str) and (file_path_or_content.startswith("http://") or file_path_or_content.startswith("https://")):
            response = requests.get(file_path_or_content)
            response.raise_for_status()
            file_content = BytesIO(response.content)
        else:
            file_content = open(file_path_or_content, 'rb')

        prs = Presentation(file_content)
        all_text = []

        # Itera sobre os slides para garantir a extração de texto de todos os slides
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extrair o texto dos shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

                # Verificar se o shape contém uma tabela
                if shape.has_table:
                    table_text = "\n[Tabela]\n"
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        table_text += row_text + "\n"
                    slide_text += table_text

            # Realizar OCR sobre imagens dentro dos slides
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image is not None:
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    ocr_text = pytesseract.image_to_string(image, lang='por')
                    slide_text += f"\n[Imagem OCR]:\n{ocr_text}\n"
            
            all_text.append(slide_text)

        # Junta todo o texto extraído dos slides
        extracted_text = "\n".join(all_text)
        return extracted_text, True

    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo PPTX: {e}")
        return "", False

    finally:
        if isinstance(file_content, BytesIO) is False and not isinstance(file_path_or_content, bytes):
            file_content.close()

# ###################   FUNCIONOU LOCALMENTE
# 

# def enhance_image(image):
#     """Melhora a imagem para OCR com nitidez e contraste."""
#     image = image.convert("L")  # Converte para escala de cinza
#     image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#     enhancer = ImageEnhance.Contrast(image)
#     return enhancer.enhance(2)  # Aumenta o contraste


# def extract_text_from_pdf(pdf_content):
#     # Verifica se o conteúdo é do tipo correto (bytes)
#     if isinstance(pdf_content, str):
#         raise ValueError("O conteúdo do PDF deve ser do tipo 'bytes' e não 'str'. Verifique a leitura do arquivo.")

#     all_text = ""

#     # Tentativa de extrair texto pesquisável e tabelas com pdfplumber
#     try:
#         with pdfplumber.open(BytesIO(pdf_content)) as pdf:
#             for page_num, page in enumerate(pdf.pages):
#                 try:
#                     # Extrai texto pesquisável
#                     page_text = page.extract_text()
#                     if page_text:
#                         all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
#                     else:
#                         print(f"Página {page_num + 1} sem texto pesquisável. Aplicando OCR...")
#                         page_image = page.to_image(resolution=300).original
#                         enhanced_image = enhance_image(page_image)
#                         ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                         all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"

#                 except Exception as e:
#                     print(f"Erro ao processar a página {page_num + 1}: {e}")

#     except Exception as e:
#         print(f"Erro ao abrir PDF com pdfplumber: {e}")

#     # Tentativa de extrair imagens e aplicar OCR com PyMuPDF
#     try:
#         doc = fitz.open(stream=BytesIO(pdf_content), filetype="pdf")
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             images = page.get_images(full=True)
#             if images:
#                 for img_index, img in enumerate(images):
#                     xref = img[0]
#                     base_image = doc.extract_image(xref)
#                     image_data = base_image["image"]
#                     image = Image.open(BytesIO(image_data))

#                     # Aplicar OCR na imagem extraída
#                     try:
#                         ocr_text = pytesseract.image_to_string(image, lang='por')
#                         all_text += f"\nImagem na página {page_num + 1}, imagem {img_index + 1}:\n{ocr_text}\n"
#                     except Exception as ocr_error:
#                         print(f"Erro ao realizar OCR na imagem da página {page_num + 1}, imagem {img_index + 1}: {ocr_error}")
#             else:
#                 print(f"Nenhuma imagem encontrada na página {page_num + 1}")

#     except Exception as e:
#         print(f"Erro ao processar imagens do PDF: {e}")

#     return all_text


# Está funcionando localmente
# # Teste da função
# if __name__ == "__main__":
#     # Exemplo de uso com leitura de arquivo local
#     pdf_path = r'D:\Repositorios\ocr-api\funcionando - V5\24_22_1300220.pdf'
#     with open(pdf_path, 'rb') as f:
#         pdf_content = f.read()  # Lendo como bytes
#     resultado = extract_text_from_pdf(pdf_content)
#     print(resultado)


# # #########################################   FUNCIONANDO MASSA FALTA APENAS IMAGEM DE PDF  ################################################

def enhance_image(image):
    """Melhora a imagem para OCR com nitidez e contraste."""
    image = image.convert("L")  # Converte para escala de cinza
    image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2)  # Aumenta o contraste
    image = ImageOps.autocontrast(image)  # Ajusta o brilho e contraste automaticamente
    return image


def download_pdf(url):
    """Baixa o PDF da URL e verifica o conteúdo."""
    try:
        response = requests.get(url, stream=True, verify=False)
        response.raise_for_status()

        # Assegura que todo o conteúdo foi baixado
        pdf_content = BytesIO(response.content)
        pdf_content.seek(0)  # Garante que o ponteiro está no início
        return pdf_content
    except Exception as e:
        print(f"Erro ao baixar o PDF: {e}")
        return None


# def extract_text_from_pdf_content(file_path):

#     """Extrai texto de PDFs pesquisáveis e não pesquisáveis, aplicando OCR."""
#     all_text = ""

#     # Verifica se o arquivo existe
#     if not os.path.exists(file_path):
#         return "", False, f"O arquivo {file_path} não existe."

#     # Primeira tentativa: extrair texto com pdfplumber
#     try:
#         with pdfplumber.open(file_path) as pdf:
#             for page_num, page in enumerate(pdf.pages):
#                 try:
#                     page_text = page.extract_text()
#                     if page_text:
#                         all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
#                     else:
#                         # Se não houver texto, aplica OCR na imagem da página
#                         page_image = page.to_image(resolution=300).original
#                         enhanced_image = enhance_image(page_image)
#                         ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                         all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"
#                 except Exception as e:
#                     print(f"Erro ao processar a página {page_num + 1}: {e}")
#     except Exception as e:
#         print(f"Erro ao abrir PDF com pdfplumber: {e}")

#     # Fallback: PyMuPDF para PDFs com imagens embutidas
#     try:
#         doc = fitz.open(file_path)
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             pix = page.get_pixmap()
#             image = Image.open(BytesIO(pix.tobytes()))
#             enhanced_image = enhance_image(image)

#             try:
#                 ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                 all_text += f"\nImagem na página {page_num + 1}:\n{ocr_text}\n"
#             except Exception as ocr_error:
#                 print(f"Erro ao realizar OCR na página {page_num + 1}: {ocr_error}")
#     except Exception as e:
#         print(f"Erro ao processar imagens do PDF: {e}")

#     return all_text, True, None


def extract_text_from_pdf_content(file_path, passwords=None):
    """
    Extrai texto de PDFs pesquisáveis e não pesquisáveis, aplicando OCR.
    Tenta abrir PDFs protegidos usando uma lista de senhas.
    """
    if passwords is None:
        passwords = []

    all_text = ""
    pdf_opened = False

    # Verifica se o arquivo existe
    if not os.path.exists(file_path):
        return "", False, f"O arquivo {file_path} não existe."

    # Tenta abrir o PDF com cada senha
    for password in passwords + [None]:  # Tenta com as senhas e sem senha
        try:
            with pdfplumber.open(file_path, password=password) as pdf:
                pdf_opened = True
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
                        else:
                            # Se não houver texto, aplica OCR na imagem da página
                            page_image = page.to_image(resolution=300).original
                            enhanced_image = enhance_image(page_image)
                            ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                            all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"
                    except Exception as e:
                        print(f"Erro ao processar a página {page_num + 1}: {e}")
                break  # Sai do loop de senhas se abrir com sucesso
        except Exception as e:
            print(f"Erro ao abrir PDF com senha '{password}': {e}")
            continue  # Tenta a próxima senha

    if not pdf_opened:
        return "", False, "Não foi possível abrir o PDF com as senhas fornecidas."

    # Fallback: PyMuPDF para PDFs com imagens embutidas
    try:
        for password in passwords + [None]:
            try:
                doc = fitz.open(file_path, password=password)
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    image = Image.open(BytesIO(pix.tobytes()))
                    enhanced_image = enhance_image(image)

                    try:
                        ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                        all_text += f"\nImagem na página {page_num + 1}:\n{ocr_text}\n"
                    except Exception as ocr_error:
                        print(f"Erro ao realizar OCR na página {page_num + 1}: {ocr_error}")
                break  # Sai do loop de senhas se abrir com sucesso
            except Exception as e:
                print(f"Erro ao abrir PDF com PyMuPDF e senha '{password}': {e}")
                continue  # Tenta a próxima senha
    except Exception as e:
        print(f"Erro ao processar imagens do PDF: {e}")

    return all_text, True, None



# ######## Arquivos Zipados ############


# def extrair_arquivos_compactados(arquivo, pasta, id):
#     # Verifica se o arquivo está compactado
#     if arquivo.endswith(('.zip', '.rar', '.7z')):
#         # Descompacta o arquivo em uma pasta
#         if arquivo.endswith('.zip'):
#             with zipfile.ZipFile(arquivo, 'r') as zip_ref:
#                 zip_ref.extractall(pasta)
#         elif arquivo.endswith('.rar'):
#             with rarfile.RarFile(arquivo, 'r') as rar_ref:
#                 rar_ref.extractall(pasta)
#         elif arquivo.endswith('.7z'):
#             with py7zr.SevenZipFile(arquivo, 'r') as seven_zip_ref:
#                 seven_zip_ref.extractall(pasta)
        
#         # Apaga o arquivo principal
#         os.remove(arquivo)
        
#         # Salva o nome, caminho e id do arquivo
#         nome_original = arquivo
#         caminho = pasta
#         id_arquivo = id
        
#         # Percorre o diretório criado para verificar se há mais arquivos compactados
#         for root, dirs, files in os.walk(pasta):
#             for file in files:
#                 arquivo_compactado = os.path.join(root, file)
#                 extrair_arquivos_compactados(arquivo_compactado, pasta, id_arquivo)
        
#         # Submete os arquivos extraídos ao processo de OCR
#         for root, dirs, files in os.walk(pasta):
#             for file in files:
#                 arquivo_extraido = os.path.join(root, file)
#                 # Chama o método de OCR aqui
#                 ocr(arquivo_extraido)
        
#         # Apaga os arquivos extraídos
#         for root, dirs, files in os.walk(pasta):
#             for file in files:
#                 arquivo_extraido = os.path.join(root, file)
#                 os.remove(arquivo_extraido)
        
#         # Salva as informações no banco de dados
#         salvar_no_banco(nome_original, caminho, id_arquivo)


######### Fim Arquivos Zipados ########

# Chamada do método para extrair arquivos compactados






########################################  Teste PDF com Senhas  ########################################


from pdf2image import convert_from_path

# Caminho para o arquivo de senhas
arquivo_senhas = r'D:\Repositorios\ocr-api\Versao6\passwords.txt'

# Ler as senhas do arquivo
with open(arquivo_senhas, 'r', encoding='utf-8') as f:
    senhas = [linha.strip() for linha in f]

senha_encontrada = False

for senha in senhas:
    try:
        with pikepdf.open(pdf_protegido, password=senha) as pdf:
            print(f'Senha correta encontrada: {senha}')
            senha_encontrada = True
            # Salvar uma cópia do PDF sem senha
            pdf.save('arquivo_sem_senha.pdf')
            print('Uma cópia do PDF sem senha foi salva como "arquivo_sem_senha.pdf".')

            # Processo de OCR
            print('Iniciando o processo de OCR...')
            # Converter o PDF em imagens (uma imagem por página)
            pages = convert_from_path('arquivo_sem_senha.pdf')

            texto_extraido = ''

            for page_number, page_data in enumerate(pages, start=1):
                # Converter a página em imagem
                page_image = page_data.convert('RGB')
                # Aplicar OCR na imagem
                texto_pagina = pytesseract.image_to_string(page_image, lang='por')
                texto_extraido += f'\n--- Página {page_number} ---\n'
                texto_extraido += texto_pagina

            # Salvar o texto extraído em um arquivo
            with open('texto_extraido.txt', 'w', encoding='utf-8') as f_texto:
                f_texto.write(texto_extraido)

            print('Processo de OCR concluído. O texto extraído foi salvo em "texto_extraido.txt".')
            break
    except pikepdf._qpdf.PasswordError:
        print(f'Senha incorreta: {senha}')
    except Exception as e:
        print(f'Erro ao tentar abrir o PDF com a senha "{senha}": {e}')

if not senha_encontrada:
    print('Nenhuma das senhas funcionou.')

# #########################################  Fim Teste PDF com Senhas  ########################################