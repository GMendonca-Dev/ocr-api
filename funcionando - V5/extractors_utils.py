from docx import Document
import fitz  # PyMuPDF
import pytesseract
from PIL import Image, ImageEnhance, ImageOps
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import json
import requests
import subprocess
from io import BytesIO, StringIO
import os
import zipfile
from pptx import Presentation
import pandas as pd
import pdfplumber
import warnings


# Suprime os avisos do tipo UserWarning, incluindo o aviso do openpyxl
warnings.simplefilter("ignore", UserWarning)

# Caminho para o executável do Tesseract no Windows (path)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Caminhos do LibreOffice e Tesseract
libreoffice_path = r"C:\Program Files\LibreOffice\program"

os.environ["PATH"] += os.pathsep + libreoffice_path

tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_xml(file_path_or_content):
    try:
        # Verifica se é conteúdo binário ou um caminho de arquivo
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string
            file_content = file_path_or_content.decode('utf-8')
        else:
            # Lê o arquivo diretamente do caminho fornecido
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = file.read()
        
        # Analisa o conteúdo XML
        root = ET.fromstring(file_content)

        # Função recursiva para extrair texto de todos os elementos XML
        def extract_text(element):
            text = element.text.strip() if element.text else ""
            for child in element:
                text += "\n" + extract_text(child)
            return text
        
        # Extrai o texto do XML
        extracted_text = extract_text(root)
        
        return extracted_text, True, ""  # Retorna o texto extraído e sucesso

    except ET.ParseError as parse_error:
        erro_msg = f"Erro ao analisar o XML: {parse_error}"
        print(erro_msg)
        return "", False, erro_msg
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo XML: {e}"
        print(erro_msg)
        return "", False, erro_msg


# Extração de texto de arquivos ODT (LibreOffice Writer)
def extract_text_from_odt(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)

        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        text_elements = [elem.text for elem in root.iter() if elem.text is not None]
        extracted_text = "\n".join(text_elements)
        return extracted_text, True
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo ODT: {e}"
        print(erro_msg)  # Exibe o erro no console
        return "", False


# Extração de texto de arquivos ODS (LibreOffice Calc)
def extract_text_from_ods(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)
        
        # Abre o arquivo 'content.xml' dentro do ODS
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Inicializa uma lista para armazenar o conteúdo extraído das células
        cell_texts = []

        # Percorre os elementos de célula no arquivo ODS
        for cell in root.iter('{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-cell'):
            # Verifica o conteúdo do texto dentro da célula
            text_content = ''.join(cell.itertext())
            if text_content:  # Adiciona o conteúdo apenas se não for vazio
                cell_texts.append(text_content)

        # Junta o conteúdo de todas as células em uma única string
        extracted_text = "\n".join(cell_texts)
        return extracted_text, True
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo ODS: {e}")
        return "", False


# Extração de texto de arquivos ODP (LibreOffice Impress)
def extract_text_from_odp(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)

        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        text_elements = [elem.text for elem in root.iter() if elem.text is not None]
        extracted_text = "\n".join(text_elements)
        return extracted_text, True
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo ODP: {e}"
        print(erro_msg)  # Exibe o erro no console
        return "", False


# Extração de texto de arquivos ODG (LibreOffice Draw)
def extract_text_from_odg(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)
        
        # Abre o arquivo 'content.xml' dentro do ODG
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Inicializa uma lista para armazenar o conteúdo extraído
        extracted_texts = []

        # Itera sobre todos os elementos de texto relevantes no arquivo ODG, incluindo texto dentro de tags aninhadas
        for elem in root.iter():
            # Verifica se o elemento possui texto (ou parte de texto) e o adiciona à lista
            if elem.text:
                extracted_texts.append(elem.text)
            if elem.tail:  # Também captura texto após a tag de fechamento
                extracted_texts.append(elem.tail)

        # Junta o conteúdo extraído em uma única string
        extracted_text = "\n".join(extracted_texts)
        return extracted_text, True
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo ODG: {e}")
        return "", False


def extract_text_from_txt(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            resultado = file_path_or_content.decode('latin1')
        else:
            with open(file_path_or_content, 'r', encoding='latin1') as file:
                resultado = file.read()
        return resultado, True
    except Exception as e:
        print(f"Erro ao ler arquivo txt: {e}")
        return "", False


# def extract_text_from_pdf(pdf_content):
    
#     try:
#         doc = fitz.open(stream=pdf_content, filetype="pdf")
#         all_text = ""
#         # ocr_status = True
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             text = page.get_text("text")
#             all_text += text

#             images = page.get_images(full=True)
#             for img in images:
#                 xref = img[0]
#                 base_image = doc.extract_image(xref)
#                 image_data = base_image["image"]
#                 image = Image.open(BytesIO(image_data))
#                 try:
#                     ocr_text = pytesseract.image_to_string(image)
#                     all_text += ocr_text
#                 except Exception as ocr_error:
#                     print(f"Erro ao realizar OCR na imagem da página {page_num}: {ocr_error}")
#                     # ocr_status = False
#         return all_text #, ocr_status
#     except Exception as e:
#         # print(f"Erro ao extrair texto do PDF: {e}")
#         return ""


def extract_text_from_xlsx(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            file_content = BytesIO(file_path_or_content)
            df = pd.read_excel(file_content, engine='openpyxl')
        else:
            df = pd.read_excel(file_path_or_content, engine='openpyxl')
        return df.to_string(), True
    except Exception as e:
        print(f"Erro ao ler arquivo xlsx: {e}")
        return "", False


def extract_text_from_csv(file_path_or_content):
    try:
        # Verifica se o conteúdo é binário (baixado de uma URL e não de um arquivo local)
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string usando StringIO
            file_content = StringIO(file_path_or_content.decode('utf-8'))
        else:
            # Se for um caminho local, abre o arquivo e lê o conteúdo
            file_content = file_path_or_content

        # Usa pandas para ler o CSV e convertê-lo em uma string
        df = pd.read_csv(file_content, on_bad_lines='skip')
        return df.to_string(), True, ""  # Retorna o conteúdo como string, status de sucesso e erro vazio
    except Exception as e:
        erro_msg = f"Erro ao ler arquivo CSV: {e}"
        print(erro_msg)
        return "", False, erro_msg


def extract_text_and_images_from_docx(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            doc = Document(BytesIO(file_path_or_content))
        else:
            doc = Document(file_path_or_content)

        all_text = ""
        extracted_text_from_images = []
        extracted_tables = []

        for para in doc.paragraphs:
            all_text += para.text + "\n"

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            extracted_tables.append(table_data)

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                text_from_image = extract_text_from_image(image_data)
                if isinstance(text_from_image, str):
                    extracted_text_from_images.append(text_from_image)

        combined_text = all_text + "\n".join(extracted_text_from_images)
        combined_text += "\n\n" + "\n".join(str(table) for table in extracted_tables)

        return combined_text, True

    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo DOCX: {e}")
        return "", False


def download_and_convert_doc_to_docx(url, temp_dir="temp"):
    try:
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)

        response = requests.get(url, verify=False)
        response.raise_for_status()

        temp_file_path = os.path.join(temp_dir, "arquivo_temp.doc")
        with open(temp_file_path, 'wb') as temp_file:
            temp_file.write(response.content)

        docx_path = convert_doc_with_libreoffice(temp_file_path, "docx")

        if not docx_path:
            print("Erro na conversão de .doc para .docx.")
            return "", False

        text, sucesso = extract_text_and_images_from_docx(docx_path)

        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        if os.path.exists(docx_path):
            os.remove(docx_path)

        return text, sucesso

    except Exception as e:
        print(f"Erro ao baixar ou processar o arquivo .doc: {e}")
        return "", False


def convert_doc_with_libreoffice(doc_path, output_format="docx", temp_dir="temp"):
    try:
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)

        output_file_path = os.path.join(temp_dir, f"arquivo_temp.{output_format}")
        command = ['soffice', '--headless', '--convert-to', output_format, doc_path, '--outdir', temp_dir]
        subprocess.run(command, check=True)

        if os.path.exists(output_file_path):
            return output_file_path
        else:
            print(f"Erro ao converter o arquivo {doc_path}")
            return None
    except Exception as e:
        print(f"Erro ao converter arquivo usando LibreOffice: {e}")
        return None


def extract_text_from_image(image_path_or_content):
    try:
        # Se for uma URL, faz o download da imagem
        if isinstance(image_path_or_content, str) and image_path_or_content.startswith(('http://', 'https://')):
            response = requests.get(image_path_or_content, stream=True)
            response.raise_for_status()
            image = Image.open(BytesIO(response.content))
        elif isinstance(image_path_or_content, bytes):
            # Caso seja conteúdo binário
            image = Image.open(BytesIO(image_path_or_content))
        elif isinstance(image_path_or_content, BytesIO):
            # Caso já seja um objeto BytesIO, usa diretamente
            image = Image.open(image_path_or_content)
        else:
            # Caso contrário, trata como um arquivo local
            image = Image.open(image_path_or_content)

        # Pré-processamento da imagem
        image = image.convert('L').resize((image.width * 2, image.height * 2), resample=Image.BICUBIC)
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(1.5)
        image = image.point(lambda x: 0 if x < 140 else 255, '1')

        # Realizar OCR na imagem
        text = pytesseract.image_to_string(image, lang='por', config='--psm 5')

        return text, True

    except Exception as e:
        print(f"Erro ao realizar OCR na imagem: {e}")
        return "", False


def extract_text_from_html(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            # Se o conteúdo for bytes (Ex: no caso de uma URL), converte para string
            file_content = file_path_or_content.decode('utf-8')
        else:
            # Se for um caminho de arquivo local
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = file.read()

        soup = BeautifulSoup(file_content, 'html.parser')
        # Extrai todo o texto visível do HTML
        text = soup.get_text(separator="\n")
        return text, True, ""  # Adicionei uma string vazia como mensagem de erro
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo HTML: {e}")
        return "", False, str(e)


def extract_text_from_json(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            file_content = json.loads(file_path_or_content.decode('utf-8'))
        else:
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = json.load(file)

        # Recorre ao conteúdo do JSON e converte tudo em uma string
        def extract_json_text(data, result=""):
            if isinstance(data, dict):
                for key, value in data.items():
                    result += f"{key}: {extract_json_text(value)}\n"
            elif isinstance(data, list):
                for item in data:
                    result += extract_json_text(item) + "\n"
            else:
                result += str(data)
            return result

        extracted_text = extract_json_text(file_content)
        return extracted_text, True, ""  # Adicionei uma string vazia como mensagem de erro
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo JSON: {e}")
        return "", False, str(e)


def extract_text_from_pptx(file_path_or_content):

    try:
        if isinstance(file_path_or_content, (bytes, BytesIO)):
            file_content = BytesIO(file_path_or_content) if isinstance(file_path_or_content, bytes) else file_path_or_content
        elif isinstance(file_path_or_content, str) and (file_path_or_content.startswith("http://") or file_path_or_content.startswith("https://")):
            response = requests.get(file_path_or_content)
            response.raise_for_status()
            file_content = BytesIO(response.content)
        else:
            file_content = open(file_path_or_content, 'rb')

        prs = Presentation(file_content)
        all_text = []

        # Itera sobre os slides para garantir a extração de texto de todos os slides
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extrair o texto dos shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

                # Verificar se o shape contém uma tabela
                if shape.has_table:
                    table_text = "\n[Tabela]\n"
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        table_text += row_text + "\n"
                    slide_text += table_text

            # Realizar OCR sobre imagens dentro dos slides
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image is not None:
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    ocr_text = pytesseract.image_to_string(image, lang='por')
                    slide_text += f"\n[Imagem OCR]:\n{ocr_text}\n"
            
            all_text.append(slide_text)

        # Junta todo o texto extraído dos slides
        extracted_text = "\n".join(all_text)
        return extracted_text, True

    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo PPTX: {e}")
        return "", False

    finally:
        if isinstance(file_content, BytesIO) is False and not isinstance(file_path_or_content, bytes):
            file_content.close()


#  ###############PDF FUNCIONANDO , SALVO IMAGENS E PDF IMAGEM
# def extract_text_from_pdf(file_path_or_url):
#     # Verifica se é uma URL ou um arquivo local
#     if file_path_or_url.startswith('http://') or file_path_or_url.startswith('https://'):
#         # Baixa o arquivo PDF da URL
#         response = requests.get(file_path_or_url, verify=False)
#         response.raise_for_status()  # Verifica se houve algum erro na requisição
#         pdf_content = response.content
#     else:
#         # Lê o arquivo PDF do diretório local
#         with open(file_path_or_url, 'rb') as file:
#             pdf_content = file.read()

#     try:
#         doc = fitz.open(stream=pdf_content, filetype="pdf")
#         all_text = ""
#         # ocr_status = True
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             text = page.get_text("text")
#             all_text += text

#             images = page.get_images(full=True)
#             for img in images:
#                 xref = img[0]
#                 base_image = doc.extract_image(xref)
#                 image_data = base_image["image"]
#                 image = Image.open(BytesIO(image_data))
#                 try:
#                     ocr_text = pytesseract.image_to_string(image)
#                     all_text += ocr_text
#                 except Exception as ocr_error:
#                     print(f"Erro ao realizar OCR na imagem da página {page_num}: {ocr_error}")
#                     # ocr_status = False
#         return all_text  # , ocr_status
    
#     except Exception as e:
#         # print(f"Erro ao extrair texto do PDF: {e}")
#         return ""



# ###################   FUNCIONOU LOCALMENTE
from PIL import ImageFilter


# def enhance_image(image):
#     """Melhora a imagem para OCR com nitidez e contraste."""
#     image = image.convert("L")  # Converte para escala de cinza
#     image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#     enhancer = ImageEnhance.Contrast(image)
#     return enhancer.enhance(2)  # Aumenta o contraste


# def extract_text_from_pdf(pdf_content):
#     # Verifica se o conteúdo é do tipo correto (bytes)
#     if isinstance(pdf_content, str):
#         raise ValueError("O conteúdo do PDF deve ser do tipo 'bytes' e não 'str'. Verifique a leitura do arquivo.")

#     all_text = ""

#     # Tentativa de extrair texto pesquisável e tabelas com pdfplumber
#     try:
#         with pdfplumber.open(BytesIO(pdf_content)) as pdf:
#             for page_num, page in enumerate(pdf.pages):
#                 try:
#                     # Extrai texto pesquisável
#                     page_text = page.extract_text()
#                     if page_text:
#                         all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
#                     else:
#                         print(f"Página {page_num + 1} sem texto pesquisável. Aplicando OCR...")
#                         page_image = page.to_image(resolution=300).original
#                         enhanced_image = enhance_image(page_image)
#                         ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                         all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"

#                 except Exception as e:
#                     print(f"Erro ao processar a página {page_num + 1}: {e}")

#     except Exception as e:
#         print(f"Erro ao abrir PDF com pdfplumber: {e}")

#     # Tentativa de extrair imagens e aplicar OCR com PyMuPDF
#     try:
#         doc = fitz.open(stream=BytesIO(pdf_content), filetype="pdf")
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             images = page.get_images(full=True)
#             if images:
#                 for img_index, img in enumerate(images):
#                     xref = img[0]
#                     base_image = doc.extract_image(xref)
#                     image_data = base_image["image"]
#                     image = Image.open(BytesIO(image_data))

#                     # Aplicar OCR na imagem extraída
#                     try:
#                         ocr_text = pytesseract.image_to_string(image, lang='por')
#                         all_text += f"\nImagem na página {page_num + 1}, imagem {img_index + 1}:\n{ocr_text}\n"
#                     except Exception as ocr_error:
#                         print(f"Erro ao realizar OCR na imagem da página {page_num + 1}, imagem {img_index + 1}: {ocr_error}")
#             else:
#                 print(f"Nenhuma imagem encontrada na página {page_num + 1}")

#     except Exception as e:
#         print(f"Erro ao processar imagens do PDF: {e}")

#     return all_text


# Está funcionando localmente
# # Teste da função
# if __name__ == "__main__":
#     # Exemplo de uso com leitura de arquivo local
#     pdf_path = r'D:\Repositorios\ocr-api\funcionando - V5\24_22_1300220.pdf'
#     with open(pdf_path, 'rb') as f:
#         pdf_content = f.read()  # Lendo como bytes
#     resultado = extract_text_from_pdf(pdf_content)
#     print(resultado)


# #########################################   FUNCIONANDO MASSA FALTA APENAS IMAGEM DE PDF  ################################################
# def enhance_image(image):
#     """Melhora a imagem para OCR com nitidez e contraste."""
#     image = image.convert("L")  # Converte para escala de cinza
#     image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#     enhancer = ImageEnhance.Contrast(image)
#     return enhancer.enhance(2)  # Aumenta o contraste


# def extract_text_from_pdf_content(pdf_content):
#     """Extrai texto de PDFs pesquisáveis e não pesquisáveis, incluindo imagens."""
#     all_text = ""

#     # Primeira tentativa: extrair texto com pdfplumber
#     try:
#         with pdfplumber.open(BytesIO(pdf_content)) as pdf:
#             for page_num, page in enumerate(pdf.pages):
#                 try:
#                     # Extrai texto pesquisável
#                     page_text = page.extract_text()
#                     if page_text:
#                         all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
#                     else:
#                         #print(f"Página {page_num + 1} sem texto pesquisável. Aplicando OCR...")
#                         # Aplica OCR se não houver texto na página
#                         page_image = page.to_image(resolution=300).original
#                         enhanced_image = enhance_image(page_image)
#                         ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                         all_text += f"{ocr_text}\n"
#                 except Exception as e:
#                     print(f"Erro ao processar a página {page_num + 1}: {e}")
#     except Exception as e:
#         print(f"Erro ao abrir PDF com pdfplumber: {e}")

#     # Fallback: Extrair imagens e aplicar OCR com PyMuPDF
#     try:
#         doc = fitz.open(stream=BytesIO(pdf_content), filetype="pdf")
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             images = page.get_images(full=True)
#             if images:
#                 for img_index, img in enumerate(images):
#                     xref = img[0]
#                     base_image = doc.extract_image(xref)
#                     image_data = base_image["image"]
#                     image = Image.open(BytesIO(image_data))

#                     # Aplicar OCR na imagem extraída
#                     try:
#                         ocr_text = pytesseract.image_to_string(image, lang='por')
#                         all_text += ocr_text
#                     except Exception as ocr_error:
#                         print(f"Erro ao realizar OCR na imagem da página {page_num + 1}, imagem {img_index + 1}: {ocr_error}")
#             else:
#                 print(f"Nenhuma imagem encontrada na página {page_num + 1}")
#     except Exception as e:
#         print(f"Erro ao processar imagens do PDF: {e}")

#     return all_text


# # #########################################   FUNCIONANDO MASSA FALTA APENAS IMAGEM DE PDF  ################################################

# 

# def extract_text_from_pdf_url(url):
#     """Combina o download e a extração de texto do PDF a partir de uma URL."""
#     pdf_content = download_pdf_from_url(url)  # Baixa o PDF da URL
#     return extract_text_from_pdf_content(pdf_content)  # Extrai o texto e/ou OCR

# url = "https://10.100.77.20/xpertis/App/Lib/arquivos/documentos/oficiosoperadoras/24_22_13022020130620_oficio02920timop00220.pdf"
# resultado = extract_text_from_pdf_url(url)
# print(resultado)




#   ################ FUNCIONANDO LOCALMENTE

# def enhance_image(image):
#     """Melhora a imagem para OCR com nitidez e contraste."""
#     image = image.convert("L")  # Converte para escala de cinza
#     image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#     enhancer = ImageEnhance.Contrast(image)
#     return enhancer.enhance(2)  # Aumenta o contraste


# def download_pdf(url):
#     """Baixa o PDF da URL e verifica o conteúdo."""
#     try:
#         response = requests.get(url, stream=True, verify=False)
#         response.raise_for_status()

#         # Assegura que todo o conteúdo foi baixado
#         pdf_content = BytesIO(response.content)
#         pdf_content.seek(0)  # Garante que o ponteiro está no início
#         return pdf_content
#     except Exception as e:
#         print(f"Erro ao baixar o PDF: {e}")
#         return None


# def extract_text_from_pdf_content(pdf_content):
#     """Extrai texto de PDFs pesquisáveis e não pesquisáveis, aplicando OCR."""
#     all_text = ""

#     try:
#         # Primeira tentativa: extrair texto com pdfplumber
#         with pdfplumber.open(pdf_content) as pdf:
#             for page_num, page in enumerate(pdf.pages):
#                 try:
#                     page_text = page.extract_text()
#                     if page_text:
#                         all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
#                     else:
#                         print(f"Página {page_num + 1} sem texto pesquisável. Aplicando OCR...")
#                         page_image = page.to_image(resolution=300).original
#                         enhanced_image = enhance_image(page_image)
#                         ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                         all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"
#                 except Exception as e:
#                     print(f"Erro ao processar a página {page_num + 1}: {e}")

#     except Exception as e:
#         print(f"Erro ao abrir PDF com pdfplumber: {e}")

#     # Fallback: Extrair imagens e aplicar OCR com PyMuPDF
#     try:
#         doc = fitz.open(stream=pdf_content, filetype="pdf")
#         for page_num in range(doc.page_count):
#             page = doc.load_page(page_num)
#             pix = page.get_pixmap()
#             image = Image.open(BytesIO(pix.tobytes()))
#             enhanced_image = enhance_image(image)

#             try:
#                 ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
#                 all_text += f"\nImagem na página {page_num + 1}:\n{ocr_text}\n"
#             except Exception as ocr_error:
#                 print(f"Erro ao realizar OCR na página {page_num + 1}: {ocr_error}")

#     except Exception as e:
#         print(f"Erro ao processar imagens do PDF: {e}")

#     return all_text

# Exemplo de uso
# url = "https://10.100.77.20/xpertis/App/Lib/arquivos/documentos/decisao/30_7_13022020130620_decisaoproc080026911.20197qst1.pdf"


# pdf_content = download_pdf(url)

# if pdf_content:
#     texto_extraido = extract_text_from_pdf_content(pdf_content)
#     print(texto_extraido)
# else:
#     print("Não foi possível baixar ou processar o PDF.")

# ######## FUNCIONANDO url


def enhance_image(image):
    """Melhora a imagem para OCR com nitidez e contraste."""
    image = image.convert("L")  # Converte para escala de cinza
    image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2)  # Aumenta o contraste
    image = ImageOps.autocontrast(image)  # Ajusta o brilho e contraste automaticamente
    return image

def download_pdf(url):
    """Baixa o PDF da URL e verifica o conteúdo."""
    try:
        response = requests.get(url, stream=True, verify=False)
        response.raise_for_status()

        # Assegura que todo o conteúdo foi baixado
        pdf_content = BytesIO(response.content)
        pdf_content.seek(0)  # Garante que o ponteiro está no início
        return pdf_content
    except Exception as e:
        print(f"Erro ao baixar o PDF: {e}")
        return None


def extract_text_from_pdf_content(pdf_content):
    """Extrai texto de PDFs pesquisáveis e não pesquisáveis, aplicando OCR."""
    if isinstance(pdf_content, str):
        pdf_content = download_pdf(pdf_content)
        if not pdf_content:
            print("Não foi possível baixar ou processar o PDF.")
            return ""
   
    all_text = ""

    # Primeira tentativa: extrair texto com pdfplumber
    try:
        with pdfplumber.open(pdf_content) as pdf:
            for page_num, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
                    else:
                        # Se não houver texto, aplica OCR na imagem da página
                        page_image = page.to_image(resolution=300).original
                        enhanced_image = enhance_image(page_image)
                        ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                        all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"
                except Exception as e:
                    print(f"Erro ao processar a página {page_num + 1}: {e}")
    except Exception as e:
        print(f"Erro ao abrir PDF com pdfplumber: {e}")

    # Fallback: PyMuPDF para PDFs com imagens embutidas
    try:
        doc = fitz.open(stream=pdf_content, filetype="pdf")
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            image = Image.open(BytesIO(pix.tobytes()))
            enhanced_image = enhance_image(image)

            try:
                ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                all_text += f"\nImagem na página {page_num + 1}:\n{ocr_text}\n"
            except Exception as ocr_error:
                print(f"Erro ao realizar OCR na página {page_num + 1}: {ocr_error}")
    except Exception as e:
        print(f"Erro ao processar imagens do PDF: {e}")

    return all_text