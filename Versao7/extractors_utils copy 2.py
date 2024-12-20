from docx import Document
import fitz  # PyMuPDF
import pytesseract
from PIL import Image, ImageEnhance, ImageOps
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import json
import requests
from io import BytesIO, StringIO
import os
import zipfile
from pptx import Presentation
import pandas as pd
import pdfplumber
import warnings
import sys
# from lxml import etree
from PIL import ImageFilter
import subprocess
import xlrd
import csv
import cv2
import numpy as np
import chardet
import tempfile
from pdfplumber import PDF
from tempfile import NamedTemporaryFile



# ######## Conversão de arquivos .doc para .docx ########
# from docx.document import Document as Document_docx
# from win32com.client import Dispatch
# ######## Fim da Conversão de arquivos .doc para .docx ########


sys.path.insert(0, './Versao7')

# Suprime os avisos do tipo UserWarning, incluindo o aviso do openpyxl
warnings.simplefilter("ignore", UserWarning)

# Caminho para o executável do Tesseract no Windows (path)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Caminhos do LibreOffice e Tesseract
libreoffice_path = r"C:\Program Files\LibreOffice\program"

os.environ["PATH"] += os.pathsep + libreoffice_path

tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_xml(file_path_or_content):
    try:
        # Verifica se é conteúdo binário ou um caminho de arquivo
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string
            file_content = file_path_or_content.decode('utf-8')
        else:
            # Lê o arquivo diretamente do caminho fornecido
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = file.read()
        
        # Analisa o conteúdo XML
        root = ET.fromstring(file_content)

        # Função recursiva para extrair texto de todos os elementos XML
        def extract_text(element):
            text = element.text.strip() if element.text else ""
            for child in element:
                text += "\n" + extract_text(child)
            return text
        
        # Extrai o texto do XML
        extracted_text = extract_text(root)
        
        return extracted_text, True, ""  # Retorna o texto extraído e sucesso

    except ET.ParseError as parse_error:
        erro_msg = f"Erro ao analisar o XML: {parse_error}"
        print(erro_msg)
        return "", False, erro_msg
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo XML: {e}"
        print(erro_msg)
        return "", False, erro_msg


# def extract_text_from_odt(file_path_or_content):
#     try:
#         # Abre o arquivo ODT a partir de um caminho ou de um objeto de arquivo
#         if isinstance(file_path_or_content, str):
#             # file_path_or_content é um caminho de arquivo
#             zip_file = zipfile.ZipFile(file_path_or_content, 'r')
#         else:
#             # file_path_or_content é um objeto de arquivo
#             zip_file = zipfile.ZipFile(file_path_or_content)

#         # Verifica se o content.xml existe no arquivo ODT
#         if 'content.xml' in zip_file.namelist():
#             with zip_file.open('content.xml') as f:
#                 tree = etree.parse(f)
#                 # Define os namespaces utilizados
#                 namespaces = {
#                     'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
#                     'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'
#                 }
#                 # Encontra todos os elementos de parágrafo
#                 paragrafos = tree.xpath('//text:p', namespaces=namespaces)
#                 extracted_text = ''
#                 for paragrafo in paragrafos:
#                     # Extrai o texto do parágrafo, incluindo elementos de texto internos
#                     texto_paragrafo = ''.join(paragrafo.xpath('.//text()', namespaces=namespaces))
#                     extracted_text += texto_paragrafo + '\n'
#             return extracted_text, True
#         else:
#             erro_msg = "O arquivo content.xml não foi encontrado no ODT."
#             print(erro_msg)
#             return "", False
#     except Exception as e:
#         erro_msg = f"Erro ao ler ou processar o arquivo ODT: {e}"
#         print(erro_msg)  # Exibe o erro no console
#         return "", False


def extract_text_from_odt(file_path_or_content):
    """
    Extrai texto de um arquivo ODT. Se o ODT contiver texto diretamente, ele será extraído.
    Caso contrário, o ODT será convertido para PDF e o texto será extraído do PDF.

    Args:
        file_path_or_content (str): Caminho para o arquivo ODT.

    Returns:
        tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
    """
    try:
        # Verifica se o arquivo ODT existe
        if not os.path.exists(file_path_or_content):
            return "", False, "O arquivo ODT não foi encontrado."

        # Diretório temporário para saída da conversão
        temp_dir = tempfile.gettempdir()
        pdf_file_path = os.path.join(temp_dir, os.path.basename(file_path_or_content).replace(".odt", ".pdf"))

        # Converte ODT para PDF usando LibreOffice
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", temp_dir, file_path_or_content],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
        except subprocess.CalledProcessError as e:
            return "", False, f"Erro durante a conversão de ODT para PDF: {e}"

        # Verifica se o arquivo PDF foi gerado
        if not os.path.exists(pdf_file_path):
            return "", False, "Falha ao converter o arquivo ODT para PDF."

        # Extrai texto do PDF
        try:
            with PDF.open(pdf_file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            return "", False, f"Erro ao extrair texto do PDF: {e}"

        # Remove o arquivo PDF temporário
        os.remove(pdf_file_path)

        # Retorna o texto extraído
        if not text.strip():
            return "", False, "Nenhum texto foi encontrado no arquivo."
        return text.strip(), True, ""

    except Exception as e:
        return "", False, f"Erro ao processar o arquivo ODT: {e}"


# Extração de texto de arquivos ODS (LibreOffice Calc)

def extract_text_from_ods(file_path_or_content):
    """
    Extrai texto de um arquivo ODS.

    Args:
        file_path_or_content (str ou bytes): Caminho para o arquivo ODS ou conteúdo em bytes.

    Returns:
        tuple: Texto extraído (str), sucesso (bool).
    """
    try:
        # Abre o arquivo ODS como um arquivo ZIP
        if isinstance(file_path_or_content, str):
            # Caminho para o arquivo
            zip_file = zipfile.ZipFile(file_path_or_content, 'r')
        else:
            # Conteúdo em bytes
            zip_file = zipfile.ZipFile(BytesIO(file_path_or_content))

        # Abre o arquivo 'content.xml' dentro do ODS
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Define os namespaces usados no arquivo ODS
        namespaces = {
            'table': 'urn:oasis:names:tc:opendocument:xmlns:table:1.0',
            'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
        }

        # Inicializa uma lista para armazenar o conteúdo extraído
        extracted_data = []

        # Itera pelas células da tabela
        for cell in root.findall('.//table:table-cell', namespaces):
            # Verifica o conteúdo do texto dentro da célula
            cell_text = []
            for text_element in cell.findall('.//text:p', namespaces):
                if text_element.text:
                    cell_text.append(text_element.text.strip())
            
            # Adiciona o conteúdo da célula à lista, se não estiver vazio
            if cell_text:
                extracted_data.append(" ".join(cell_text))

        # Junta o conteúdo de todas as células em uma única string
        extracted_text = "\n".join(extracted_data)

        # Fecha o arquivo ZIP
        zip_file.close()

        # Verifica se algum texto foi extraído
        if not extracted_text.strip():
            return "", False

        return extracted_text, True

    except Exception as e:
        erro_msg = f"Erro ao processar o arquivo ODS: {e}"
        print(erro_msg)
        return "", False


# Extração de texto de arquivos ODP (LibreOffice Impress)
def extract_text_from_odp(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)

        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        text_elements = [elem.text for elem in root.iter() if elem.text is not None]
        extracted_text = "\n".join(text_elements)
        return extracted_text, True
    except Exception as e:
        erro_msg = f"Erro ao ler ou processar o arquivo ODP: {e}"
        print(erro_msg)  # Exibe o erro no console
        return "", False


# Extração de texto de arquivos ODG (LibreOffice Draw)
def extract_text_from_odg(file_path_or_content):
    try:
        # Garante que file_path_or_content seja tratado como um objeto de arquivo
        if isinstance(file_path_or_content, str):
            with open(file_path_or_content, 'rb') as f:
                zip_file = zipfile.ZipFile(f)
        else:
            zip_file = zipfile.ZipFile(file_path_or_content)
        
        # Abre o arquivo 'content.xml' dentro do ODG
        with zip_file.open('content.xml') as f:
            tree = ET.parse(f)
            root = tree.getroot()

        # Inicializa uma lista para armazenar o conteúdo extraído
        extracted_texts = []

        # Itera sobre todos os elementos de texto relevantes no arquivo ODG, incluindo texto dentro de tags aninhadas
        for elem in root.iter():
            # Verifica se o elemento possui texto (ou parte de texto) e o adiciona à lista
            if elem.text:
                extracted_texts.append(elem.text)
            if elem.tail:  # Também captura texto após a tag de fechamento
                extracted_texts.append(elem.tail)

        # Junta o conteúdo extraído em uma única string
        extracted_text = "\n".join(extracted_texts)
        return extracted_text, True
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo ODG: {e}")
        return "", False


# def extract_text_from_txt(file_path_or_content):
#     """
#     Extrai o texto de um arquivo de texto (formato .txt).

#     O arquivo pode ser lido como um objeto de bytes ou como um caminho de arquivo no sistema de arquivos.
#     O encoding padr o  'latin1'.
#     """
  
#     try:
#         if isinstance(file_path_or_content, bytes):
#             resultado = file_path_or_content.decode('latin1')
#         else:
#             with open(file_path_or_content, 'r', encoding='latin1') as file:
#                 resultado = file.read()
#         return resultado, True
#     except Exception as e:
#         print(f"Erro ao ler arquivo txt: {e}")
#         return "", False


def extract_text_from_txt(file_path):
    """
    Extrai texto de um arquivo TXT e corrige a codificação automaticamente.

    Args:
        file_path (str): Caminho para o arquivo TXT.

    Returns:
        tuple: Texto corrigido (str), sucesso (bool), mensagem de erro (str).
    """
    try:
        # Detecta a codificação do arquivo
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            detected_encoding = chardet.detect(raw_data)['encoding']

        # Lê o arquivo usando a codificação detectada
        with open(file_path, 'r', encoding=detected_encoding) as file:
            text = file.read()

        # Normaliza o texto removendo espaços extras
        text = text.strip()

        return text, True, ""

    except Exception as e:
        erro_msg = f"Erro ao processar o arquivo TXT: {e}"
        print(erro_msg)
        return "", False, erro_msg


def extract_text_from_odf(file_path, extension):
    """
    Extrai o texto de um arquivo ODF (OpenDocument Format).

    Parameters
    ----------
    file_path : str
        Caminho do arquivo ODF.
    extension : str
        Extens o do arquivo ODF. Pode ser 'odt', 'ods', 'odp' ou 'odg'.

    Returns
    -------
    tuple
        Tupla contendo o texto extra do, um booleano indicando se a extra o foi bem-sucedida e
        uma mensagem de erro (caso a extra o tenha falhado).

    Raises
    ------
    Exception
        Caso ocorra um erro durante a extra o de texto do arquivo ODF.
    """
    try:
        if extension == 'odt':
            # Corrigir aqui - pegar apenas o primeiro elemento da tupla
            texto, sucesso = extract_text_from_odt(file_path)
            return texto, sucesso, None
        elif extension == 'ods':
            texto, sucesso = extract_text_from_ods(file_path)
            return texto, sucesso, None
        elif extension == 'odp':
            texto, sucesso = extract_text_from_odp(file_path)
            return texto, sucesso, None
        elif extension == 'odg':
            texto, sucesso = extract_text_from_odg(file_path)
            return texto, sucesso, None
    except Exception as e:
        erro_msg = f"Erro ao processar arquivo ODF ({extension}): {e}"
        return "", False, erro_msg


# def extract_text_from_odf(file_path, extension):
#     try:
#         if extension == 'odt':
#             return extract_text_from_odt(file_path), True, None
#         elif extension == 'ods':
#             return extract_text_from_ods(file_path), True, None
#         elif extension == 'odp':
#             return extract_text_from_odp(file_path), True, None
#         elif extension == 'odg':
#             return extract_text_from_odg(file_path), True, None
#     except Exception as e:
#         erro_msg = f"Erro ao processar arquivo ODF ({extension}): {e}"
#         return "", False, erro_msg


# def extract_text_from_xlsx(file_path_or_content):

#     try:
#         if isinstance(file_path_or_content, bytes):
#             file_content = BytesIO(file_path_or_content)
#             df = pd.read_excel(file_content, engine='openpyxl')
#         else:
#             df = pd.read_excel(file_path_or_content, engine='openpyxl')
#         return df.to_string(), True
#     except Exception as e:
#         print(f"Erro ao ler arquivo xlsx: {e}")
#         return "", False


def extract_text_from_xlsx(file_path_or_content):
    """
    Extrai o texto de todas as planilhas de um arquivo .xlsx.

    Parâmetros:
    - file_path_or_content (str ou bytes): Caminho para o arquivo .xlsx ou conteúdo em bytes.

    Retorna:
    - texto_extraido (str): O texto extraído de todas as planilhas.

    Lança:
    - FileNotFoundError: Se o arquivo não existir.
    - ValueError: Se o arquivo não for um arquivo Excel válido.
    - TypeError: Se o parâmetro 'file_path_or_content' não for do tipo esperado.
    - Exception: Para outros erros inesperados.
    """
    if isinstance(file_path_or_content, bytes):
        # Conteúdo em bytes
        file_content = BytesIO(file_path_or_content)
        excel_file = file_content
    elif isinstance(file_path_or_content, str):
        # Caminho do arquivo
        if not os.path.exists(file_path_or_content):
            raise FileNotFoundError(f"O arquivo '{file_path_or_content}' não existe.")
        excel_file = file_path_or_content
    else:
        raise TypeError("O parâmetro 'file_path_or_content' deve ser um caminho de arquivo (str) ou conteúdo em bytes.")

    try:
        # Lê todas as planilhas do arquivo Excel
        df_dict = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
        texto_extraido = ''

        for nome_planilha, df in df_dict.items():
            texto_extraido += f"=== Planilha: {nome_planilha} ===\n"
            texto_extraido += df.to_string(index=False)
            texto_extraido += '\n\n'

        return texto_extraido

    except FileNotFoundError:
        # O arquivo não foi encontrado durante a leitura
        raise FileNotFoundError(f"O arquivo '{file_path_or_content}' não foi encontrado.")
    except ValueError as ve:
        # Erro ao ler o arquivo Excel
        raise ValueError(f"Erro ao ler o arquivo xlsx: {ve}")
    except Exception as e:
        # Outros erros inesperados
        raise Exception(f"Erro inesperado ao processar o arquivo: {e}")


def extract_text_from_xls(file_path):
  
    """
    Detecta o formato real do arquivo e extrai o texto.

    Parâmetros:
    caminho_arquivo_xls (str): Caminho para o arquivo de entrada.

    Retorna:
    str: Texto extraído do documento.
    """
    # Verifica o formato real do arquivo
    with open(file_path, 'rb') as f:
        inicio = f.read(1024).lower()

    try:
        if b'<html' in inicio or b'<table' in inicio:
            # Trata como HTML
            texto = extract_text_from_html(file_path)
        elif inicio.startswith(b'pk'):
            # Trata como XLSX (arquivo zip)
            texto = extract_text_from_xlsx(file_path)
        else:
            try:
                workbook = xlrd.open_workbook(file_path)
                texto = ''
                for sheet in workbook.sheets():
                    texto += f"=== Planilha: {sheet.name} ===\n"
                    for row_idx in range(sheet.nrows):
                        row = sheet.row_values(row_idx)
                        texto += ', '.join(map(str, row)) + '\n'
                # return texto
            except Exception as e:
                print(f"Erro ao extrair texto do XLS: {e}")
                raise e
        return texto
    except Exception as e:
        print(f"Erro ao extrair texto: {e}")
        raise e


def extract_text_from_xltx(file_path_or_content):
    """
    Extrai o texto de todas as planilhas de um arquivo .xltx.

    Args:
    - file_path_or_content (str ou bytes): Caminho para o arquivo .xltx ou conteúdo em bytes.

    Retorna:
    - texto_extraido (str): O texto extraído de todas as planilhas.
    """
    if isinstance(file_path_or_content, bytes):
        # Conteúdo em bytes
        file_content = BytesIO(file_path_or_content)
        excel_file = file_content
    elif isinstance(file_path_or_content, str):
        # Caminho do arquivo
        if not os.path.exists(file_path_or_content):
            raise FileNotFoundError(f"O arquivo '{file_path_or_content}' não existe.")
        excel_file = file_path_or_content
    else:
        raise TypeError("O parâmetro 'file_path_or_content' deve ser um caminho de arquivo (str) ou conteúdo em bytes.")

    try:
        # Lê todas as planilhas do arquivo Excel
        df_dict = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
        texto_extraido = ''

        for nome_planilha, df in df_dict.items():
            # Substitui os valores NaN por strings vazias
            df = df.fillna('')
            texto_extraido += f"=== Planilha: {nome_planilha} ===\n"
            texto_extraido += df.to_string(index=False)
            texto_extraido += '\n\n'

        return texto_extraido

    except FileNotFoundError:
        raise FileNotFoundError(f"O arquivo '{file_path_or_content}' não foi encontrado.")
    except ValueError as ve:
        raise ValueError(f"Erro ao ler o arquivo xltx: {ve}")
    except Exception as e:
        raise Exception(f"Erro inesperado ao processar o arquivo: {e}")


# def extract_text_from_csv(file_path_or_content):
#     """
#     Extrai o texto de um arquivo CSV.

#     Parâmetros:
#     - file_path_or_content (str ou bytes): Caminho para o arquivo CSV ou conteúdo em bytes.

#     Retorna:
#     - texto_extraido (str): O texto extraído do arquivo CSV.
#     - sucesso (bool): Indica se a extração foi bem-sucedida.
#     - erro_msg (str): Mensagem de erro, se houver.

#     Exceções:
#     - Gera e imprime uma mensagem de erro se ocorrer qualquer exceção durante a leitura do arquivo CSV.
#     """
#     try:
#         # Verifica se o conteúdo é binário (baixado de uma URL e não de um arquivo local)
#         if isinstance(file_path_or_content, bytes):
#             # Converte bytes para string usando StringIO
#             file_content = StringIO(file_path_or_content.decode('utf-8'))
#         else:
#             # Se for um caminho local, abre o arquivo e lê o conteúdo
#             file_content = file_path_or_content

#         # Usa pandas para ler o CSV e convertê-lo em uma string
#         df = pd.read_csv(file_content, on_bad_lines='skip')
#         return df.to_string(), True, ""  # Retorna o conteúdo como string, status de sucesso e erro vazio
#     except Exception as e:
#         erro_msg = f"Erro ao ler arquivo CSV: {e}"
#         print(erro_msg)
#         return "", False, erro_msg


# def extract_csv_with_libreoffice(file_path, output_dir=None):


def extract_text_from_csv(file_path_or_content, delimiter=None):
    """
    Extrai o texto de um arquivo CSV.

    Parâmetros:
    - file_path_or_content (str ou bytes): Caminho para o arquivo CSV ou conteúdo em bytes.
    - delimiter (str, opcional): Separador usado no arquivo CSV. Se não fornecido, será detectado automaticamente.

    Retorna:
    - texto_extraido (str): O texto extraído do arquivo CSV.
    - sucesso (bool): Indica se a extração foi bem-sucedida.
    - erro_msg (str): Mensagem de erro, se houver.
    """
    # Lista de codificações para tentar
    encodings = ['utf-8', 'latin1', 'cp1252']

    try:
        # Verifica se o conteúdo é binário (baixado de uma URL ou bytes)
        if isinstance(file_path_or_content, bytes):
            # Converte bytes para string usando StringIO
            file_content = StringIO(file_path_or_content.decode('utf-8'))
        else:
            # Verifica se o caminho existe
            if not os.path.exists(file_path_or_content):
                raise FileNotFoundError(f"O arquivo '{file_path_or_content}' não foi encontrado.")
            file_content = file_path_or_content

        # Detecta automaticamente o delimitador, se não fornecido
        if delimiter is None:
            with open(file_path_or_content, 'r', encoding='cp1252', errors='ignore') as f:
                sample = f.read(1024)
                try:
                    delimiter = csv.Sniffer().sniff(sample).delimiter
                except csv.Error:
                    delimiter = ','  # Fallback para vírgula como padrão

        # Tenta ler o arquivo com diferentes codificações
        for encoding in encodings:
            try:
                df = pd.read_csv(
                    file_content,
                    on_bad_lines='skip',  # Ignora linhas problemáticas
                    dtype=str,  # Força os dados como strings
                    quotechar="'",  # Trata aspas
                    quoting=csv.QUOTE_MINIMAL,  # Aspas mínimas
                    escapechar='\\',  # Trata caracteres de escape
                    delimiter=delimiter,  # Usa o delimitador detectado ou especificado
                    encoding=encoding  # Tenta com a codificação atual
                )
                break  # Sai do loop se a leitura for bem-sucedida
            except UnicodeDecodeError:
                continue  # Tenta a próxima codificação
        else:
            # Se nenhuma codificação funcionar
            raise UnicodeDecodeError(f"Não foi possível decodificar o arquivo usando as codificações: {encodings}")

        if df.empty:
            return "", False, f"O arquivo CSV está vazio ou não contém dados legíveis. Separador detectado: '{delimiter}'."

        # Converte o DataFrame em uma string formatada
        texto_extraido = df.to_string(index=False, na_rep="")
        return texto_extraido, True, ""

    except FileNotFoundError as fnf_error:
        erro_msg = str(fnf_error)
        print(erro_msg)
        return "", False, erro_msg
    except pd.errors.EmptyDataError:
        erro_msg = f"O arquivo CSV está vazio ou mal formatado. Separador usado: '{delimiter}'."
        print(erro_msg)
        return "", False, erro_msg
    except Exception as e:
        erro_msg = f"Erro ao ler arquivo CSV: {e}"
        print(erro_msg)
        return "", False, erro_msg
    

# def extract_text_and_images_from_docx(file_path_or_content):
#     try:
#         if isinstance(file_path_or_content, bytes):
#             doc = Document(BytesIO(file_path_or_content))
#         else:
#             doc = Document(file_path_or_content)

#         all_text = ""
#         extracted_text_from_images = []
#         extracted_tables = []

#         for para in doc.paragraphs:
#             all_text += para.text + "\n"

#         for table in doc.tables:
#             table_data = []
#             for row in table.rows:
#                 row_data = [cell.text for cell in row.cells]
#                 table_data.append(row_data)
#             extracted_tables.append(table_data)

#         for rel in doc.part.rels.values():
#             if "image" in rel.target_ref:
#                 image_data = rel.target_part.blob
#                 text_from_image = extract_text_from_image(image_data)
#                 if isinstance(text_from_image, str):
#                     extracted_text_from_images.append(text_from_image)

#         combined_text = all_text + "\n".join(extracted_text_from_images)
#         combined_text += "\n\n" + "\n".join(str(table) for table in extracted_tables)

#         return combined_text, True

#     except Exception as e:
#         print(f"Erro ao ler ou processar o arquivo DOCX: {e}")
#         return "", False


# def extract_text_and_images_from_docx(file_path_or_content):
#     """
#     Extrai texto e imagens (com OCR) de arquivos DOCX.

#     Args:
#         file_path_or_content (str ou bytes): Caminho para o arquivo DOCX ou conteúdo em bytes.

#     Returns:
#         tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
#     """
#     try:
#         if isinstance(file_path_or_content, bytes):
#             doc = Document(BytesIO(file_path_or_content))
#         else:
#             doc = Document(file_path_or_content)

#         all_text = ""
#         extracted_text_from_images = []

#         # Extrai o texto de parágrafos
#         for para in doc.paragraphs:
#             all_text += para.text + "\n"

#         # Extrai o texto de tabelas
#         for table in doc.tables:
#             for row in table.rows:
#                 row_text = "\t".join(cell.text.strip() for cell in row.cells)
#                 all_text += row_text + "\n"

#         # Processa imagens embutidas no documento
#         for rel in doc.part.rels.values():
#             if "image" in rel.target_ref:
#                 image_data = rel.target_part.blob
#                 image = Image.open(BytesIO(image_data))

#                 # Pré-processamento da imagem
#                 image = image.convert("L")  # Converte para escala de cinza
#                 image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#                 enhancer = ImageEnhance.Contrast(image)
#                 image = enhancer.enhance(2)  # Aumenta o contraste
#                 image = ImageOps.autocontrast(image)  # Ajusta brilho e contraste automaticamente

#                 # Extrai texto da imagem com OCR
#                 text_from_image = pytesseract.image_to_string(image, lang='por', config='--psm 6')
#                 if text_from_image.strip():
#                     extracted_text_from_images.append(text_from_image)

#         # Combina texto extraído de parágrafos, tabelas e imagens
#         combined_text = all_text.strip()
#         if extracted_text_from_images:
#             combined_text += "\n\n--- Texto das Imagens ---\n" + "\n".join(extracted_text_from_images)

#         return combined_text, True, ""

#     except Exception as e:
#         erro_msg = f"Erro ao processar o arquivo DOCX: {e}"
#         print(erro_msg)
#         return "", False, erro_msg


# def extract_text_and_images_from_docx(file_path_or_content):
#     """
#     Extrai texto e realiza OCR em imagens embutidas em arquivos DOCX.

#     Args:
#         file_path_or_content (str ou bytes): Caminho para o arquivo DOCX ou conteúdo em bytes.

#     Returns:
#         tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
#     """
#     try:
#         if isinstance(file_path_or_content, bytes):
#             doc = Document(BytesIO(file_path_or_content))
#         else:
#             doc = Document(file_path_or_content)

#         all_text = ""
#         extracted_text_from_images = []

#         # Extrai texto de parágrafos
#         for para in doc.paragraphs:
#             all_text += para.text + "\n"

#         # Extrai texto de tabelas
#         for table in doc.tables:
#             for row in table.rows:
#                 row_text = "\t".join(cell.text.strip() for cell in row.cells)
#                 all_text += row_text + "\n"

#         # Processa imagens embutidas no documento
#         for rel in doc.part.rels.values():
#             if "image" in rel.target_ref:
#                 image_data = rel.target_part.blob
#                 image = Image.open(BytesIO(image_data))

#                 # Pré-processamento da imagem
#                 image = image.convert("L")  # Converte para escala de cinza
#                 image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#                 enhancer = ImageEnhance.Contrast(image)
#                 image = enhancer.enhance(2)  # Aumenta o contraste
#                 image = ImageOps.autocontrast(image)  # Ajusta brilho e contraste automaticamente

#                 # Extrai texto da imagem com OCR
#                 text_from_image = pytesseract.image_to_string(image, lang='por', config='--psm 6')
#                 if text_from_image.strip():
#                     extracted_text_from_images.append(text_from_image)

#         # Combina texto extraído de parágrafos, tabelas e imagens
#         combined_text = all_text.strip()
#         if extracted_text_from_images:
#             combined_text += "\n\n--- Texto das Imagens ---\n" + "\n".join(extracted_text_from_images)

#         return combined_text, True, ""

#     except Exception as e:
#         erro_msg = f"Erro ao processar o arquivo DOCX: {e}"
#         print(erro_msg)
#         return "", False, erro_msg


def extract_text_and_images_from_docx(file_path_or_content):
    """
    Extrai texto e realiza OCR em imagens embutidas em arquivos DOCX.

    Args:
        file_path_or_content (str ou bytes): Caminho para o arquivo DOCX ou conteúdo em bytes.

    Returns:
        tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
    """
    try:
        if isinstance(file_path_or_content, bytes):
            doc = Document(BytesIO(file_path_or_content))
        else:
            doc = Document(file_path_or_content)

        all_text = ""
        extracted_text_from_images = []

        # Extrai texto de parágrafos
        for para in doc.paragraphs:
            all_text += para.text + "\n"

        # Extrai texto de tabelas
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                all_text += row_text + "\n"

        # Processa imagens embutidas no documento usando a função de extração de imagens
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                with BytesIO(image_data) as image_stream:
                    # Salva a imagem em memória e chama a função de extração
                    image_path = image_stream  # Simula o caminho do arquivo
                    text_from_image, success, error_msg = extract_text_from_image(image_path)

                    if success and text_from_image.strip():
                        extracted_text_from_images.append(text_from_image)

        # Combina texto extraído de parágrafos, tabelas e imagens
        combined_text = all_text.strip()
        if extracted_text_from_images:
            combined_text += "\n\n--- Texto das Imagens ---\n" + "\n".join(extracted_text_from_images)

        return combined_text, True, ""

    except Exception as e:
        erro_msg = f"Erro ao processar o arquivo DOCX: {e}"
        print(erro_msg)
        return "", False, erro_msg


def download_and_convert_doc_to_docx(file_path, diretorio_saida=None):
    """
    Converte um arquivo .doc para .docx usando LibreOffice.

    Parâmetros:
    caminho_arquivo_doc (str): Caminho para o arquivo .doc de entrada.
    diretorio_saida (str, opcional): Diretório onde o arquivo convertido será salvo. 
                                     Se None, usa o diretório do arquivo original.

    Retorna:
    str: Caminho para o arquivo .docx convertido.
    """
    if diretorio_saida is None:
        diretorio_saida = os.path.dirname(file_path)
    
    try:
        # Executa o comando de conversão
        subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'docx',
            file_path,
            '--outdir', diretorio_saida
        ], check=True)
        
        # Obtém o nome base do arquivo sem extensão
        nome_base = os.path.splitext(os.path.basename(file_path))[0]
        caminho_docx = os.path.join(diretorio_saida, f"{nome_base}.docx")
        
        if os.path.exists(caminho_docx):
            return caminho_docx
        else:
            raise FileNotFoundError(f"Arquivo convertido {caminho_docx} não encontrado.")
    
    except subprocess.CalledProcessError as e:
        print(f"Erro na conversão do arquivo: {e}")
        raise e


# def extract_text_from_image(image_path):
#     """
#     Extrai texto de uma imagem (JPEG ou PNG) utilizando OCR, com pré-processamento.

#     Args:
#         image_path (str): Caminho para o arquivo de imagem.

#     Returns:
#         tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
#     """
#     try:
#         # Abre a imagem
#         image = Image.open(image_path)

#         # Pré-processamento da imagem
#         image = image.convert("L")  # Converte para escala de cinza
#         image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
#         enhancer = ImageEnhance.Contrast(image)
#         image = enhancer.enhance(2)  # Aumenta o contraste
#         image = ImageOps.autocontrast(image)  # Ajusta brilho e contraste automaticamente

#         # Aplica OCR na imagem
#         text = pytesseract.image_to_string(image, lang='por', config='--psm 6')

#         # Remove espaços extras e normaliza o texto
#         text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])

#         return text, True, ""

#     except Exception as e:
#         erro_msg = f"Erro ao realizar OCR na imagem: {e}"
#         print(erro_msg)
#         return "", False, erro_msg


def extract_text_from_image(image_path):
    """
    Extrai texto completo de uma imagem (JPEG ou PNG), incluindo tabelas, utilizando OCR.

    Args:
        image_path (str): Caminho para o arquivo de imagem.

    Returns:
        tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
    """
    try:
        # Abre a imagem com Pillow
        image = Image.open(image_path)

        # Pré-processamento básico (escalar de cinza, contraste)
        image = image.convert("L")  # Converte para escala de cinza
        image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2)  # Aumenta o contraste
        image = ImageOps.autocontrast(image)  # Ajusta brilho e contraste automaticamente

        # Aplica OCR para texto geral
        text_general = pytesseract.image_to_string(image, lang='por', config='--psm 6')

        # Remove espaços extras e normaliza o texto geral
        text_general = "\n".join([line.strip() for line in text_general.splitlines() if line.strip()])

        # Converte a imagem para formato OpenCV
        image_cv = cv2.cvtColor(np.array(image), cv2.COLOR_GRAY2BGR)

        # Detecção de tabelas com OpenCV
        gray = cv2.cvtColor(image_cv, cv2.COLOR_BGR2GRAY)
        blurred = cv2.GaussianBlur(gray, (5, 5), 0)
        edges = cv2.Canny(blurred, 50, 150)

        # Detecta contornos
        contours, _ = cv2.findContours(edges, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

        # Filtra contornos que podem representar células de tabela
        table_contours = [c for c in contours if cv2.contourArea(c) > 100]

        text_table = ""
        if table_contours:
            # Pré-processa novamente para OCR de tabela
            custom_config = '--psm 11'  # Modo de layout para tabelas
            text_table = pytesseract.image_to_string(image, lang='por', config=custom_config)

            # Remove espaços extras e normaliza o texto das tabelas
            text_table = "\n".join([line.strip() for line in text_table.splitlines() if line.strip()])

        # Combina texto geral e texto de tabelas
        combined_text = text_general
        if text_table:
            combined_text += "\n\n--- Texto de Tabelas ---\n" + text_table

        return combined_text, True, ""

    except Exception as e:
        erro_msg = f"Erro ao realizar OCR na imagem: {e}"
        print(erro_msg)
        return "", False, erro_msg


# def extract_text_from_html(file_path_or_content):
#     try:
#         if isinstance(file_path_or_content, bytes):
#             # Se o conteúdo for bytes (Ex: no caso de uma URL), converte para string
#             file_content = file_path_or_content.decode('utf-8')
#         else:
#             # Se for um caminho de arquivo local
#             with open(file_path_or_content, 'r', encoding='utf-8') as file:
#                 file_content = file.read()

#         soup = BeautifulSoup(file_content, 'html.parser')
#         # Extrai todo o texto visível do HTML
#         text = soup.get_text(separator="\n")
#         return text, True, ""  # Adicionei uma string vazia como mensagem de erro
#     except Exception as e:
#         print(f"Erro ao ler ou processar o arquivo HTML: {e}")
#         return "", False, str(e)


# def extract_text_from_html(file_path_or_content):
#     """
#     Extrai texto visível de um arquivo HTML ou conteúdo HTML.

#     Args:
#         file_path_or_content (str ou bytes): Caminho para o arquivo HTML ou conteúdo HTML em bytes.

#     Returns:
#         tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
#     """
#     try:
#         # Lê o conteúdo HTML
#         if isinstance(file_path_or_content, bytes):
#             # Caso seja conteúdo em bytes (ex: de uma URL), converte para string
#             file_content = file_path_or_content.decode('utf-8', errors='replace')
#         elif isinstance(file_path_or_content, str):
#             if "<html" in file_path_or_content.lower():
#                 # Caso seja diretamente o conteúdo HTML em string
#                 file_content = file_path_or_content
#             else:
#                 # Caso seja um caminho de arquivo local
#                 with open(file_path_or_content, 'r', encoding='utf-8', errors='replace') as file:
#                     file_content = file.read()
#         else:
#             raise TypeError("O parâmetro 'file_path_or_content' deve ser str ou bytes.")

#         # Usa BeautifulSoup para processar o HTML
#         soup = BeautifulSoup(file_content, 'html.parser')

#         # Remove scripts, estilos e outros elementos não visíveis
#         for element in soup(["script", "style", "noscript", "meta", "link"]):
#             element.extract()

#         # Extrai texto visível do HTML
#         text = soup.get_text(separator="\n")
#         text = "\n".join(line.strip() for line in text.splitlines() if line.strip())  # Remove linhas vazias/extras

#         return text, True, ""
#     except Exception as e:
#         erro_msg = f"Erro ao processar o arquivo HTML: {e}"
#         print(erro_msg)
#         return "", False, erro_msg


def extract_text_from_html(file_path_or_content):
    """
    Extrai texto visível de um arquivo HTML ou conteúdo HTML, corrigindo problemas de codificação.

    Args:
        file_path_or_content (str ou bytes): Caminho para o arquivo HTML ou conteúdo HTML em bytes.

    Returns:
        tuple: Texto extraído (str), sucesso (bool), mensagem de erro (str).
    """
    try:
        # Detecta codificação e lê o conteúdo HTML
        if isinstance(file_path_or_content, bytes):
            # Detecta a codificação se o conteúdo for em bytes
            detected_encoding = chardet.detect(file_path_or_content)['encoding']
            file_content = file_path_or_content.decode(detected_encoding, errors='replace')
        elif isinstance(file_path_or_content, str):
            if "<html" in file_path_or_content.lower():
                # Caso seja diretamente o conteúdo HTML em string
                file_content = file_path_or_content
            else:
                # Caso seja um caminho de arquivo local
                with open(file_path_or_content, 'rb') as file:
                    raw_data = file.read()
                    detected_encoding = chardet.detect(raw_data)['encoding']
                    file_content = raw_data.decode(detected_encoding, errors='replace')
        else:
            raise TypeError("O parâmetro 'file_path_or_content' deve ser str ou bytes.")

        # Usa BeautifulSoup para processar o HTML
        soup = BeautifulSoup(file_content, 'html.parser')

        # Remove scripts, estilos e outros elementos não visíveis
        for element in soup(["script", "style", "noscript", "meta", "link"]):
            element.extract()

        # Extrai texto visível do HTML
        text = soup.get_text(separator="\n")
        text = "\n".join(line.strip() for line in text.splitlines() if line.strip())  # Remove linhas vazias/extras

        return text, True, ""
    except Exception as e:
        erro_msg = f"Erro ao processar o arquivo HTML: {e}"
        print(erro_msg)
        return "", False, erro_msg


def extract_text_from_json(file_path_or_content):
    try:
        if isinstance(file_path_or_content, bytes):
            file_content = json.loads(file_path_or_content.decode('utf-8'))
        else:
            with open(file_path_or_content, 'r', encoding='utf-8') as file:
                file_content = json.load(file)

        # Recorre ao conteúdo do JSON e converte tudo em uma string
        def extract_json_text(data, result=""):
            if isinstance(data, dict):
                for key, value in data.items():
                    result += f"{key}: {extract_json_text(value)}\n"
            elif isinstance(data, list):
                for item in data:
                    result += extract_json_text(item) + "\n"
            else:
                result += str(data)
            return result

        extracted_text = extract_json_text(file_content)
        return extracted_text, True, ""  # Adicionei uma string vazia como mensagem de erro
    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo JSON: {e}")
        return "", False, str(e)


def extract_text_from_pptx(file_path_or_content):

    try:
        if isinstance(file_path_or_content, (bytes, BytesIO)):
            file_content = BytesIO(file_path_or_content) if isinstance(file_path_or_content, bytes) else file_path_or_content
        elif isinstance(file_path_or_content, str) and (file_path_or_content.startswith("http://") or file_path_or_content.startswith("https://")):
            response = requests.get(file_path_or_content)
            response.raise_for_status()
            file_content = BytesIO(response.content)
        else:
            file_content = open(file_path_or_content, 'rb')

        prs = Presentation(file_content)
        all_text = []

        # Itera sobre os slides para garantir a extração de texto de todos os slides
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            # Extrair o texto dos shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

                # Verificar se o shape contém uma tabela
                if shape.has_table:
                    table_text = "\n[Tabela]\n"
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        table_text += row_text + "\n"
                    slide_text += table_text

            # Realizar OCR sobre imagens dentro dos slides
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image is not None:
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    ocr_text = pytesseract.image_to_string(image, lang='por')
                    slide_text += f"\n[Imagem OCR]:\n{ocr_text}\n"
            
            all_text.append(slide_text)

        # Junta todo o texto extraído dos slides
        extracted_text = "\n".join(all_text)
        return extracted_text, True

    except Exception as e:
        print(f"Erro ao ler ou processar o arquivo PPTX: {e}")
        return "", False

    finally:
        if isinstance(file_content, BytesIO) is False and not isinstance(file_path_or_content, bytes):
            file_content.close()


def enhance_image(image):
    """Melhora a imagem para OCR com nitidez e contraste."""
    image = image.convert("L")  # Converte para escala de cinza
    image = image.filter(ImageFilter.SHARPEN)  # Aumenta a nitidez
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2)  # Aumenta o contraste
    image = ImageOps.autocontrast(image)  # Ajusta o brilho e contraste automaticamente
    return image


def download_pdf(url):
    """Baixa o PDF da URL e verifica o conteúdo."""
    try:
        response = requests.get(url, stream=True, verify=False)
        response.raise_for_status()

        # Assegura que todo o conteúdo foi baixado
        pdf_content = BytesIO(response.content)
        pdf_content.seek(0)  # Garante que o ponteiro está no início
        return pdf_content
    except Exception as e:
        print(f"Erro ao baixar o PDF: {e}")
        return None


def extract_text_from_pdf_content(file_path):

    """Extrai texto de PDFs pesquisáveis e não pesquisáveis, aplicando OCR."""
    all_text = ""

    # Verifica se o arquivo existe
    if not os.path.exists(file_path):
        return "", False, f"O arquivo {file_path} não existe."

    # Primeira tentativa: extrair texto com pdfplumber
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        all_text += f"\nPágina {page_num + 1}:\n{page_text}\n"
                    else:
                        # Se não houver texto, aplica OCR na imagem da página
                        page_image = page.to_image(resolution=300).original
                        enhanced_image = enhance_image(page_image)
                        ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                        all_text += f"\nPágina {page_num + 1} (OCR):\n{ocr_text}\n"
                except Exception as e:
                    print(f"Erro ao processar a página {page_num + 1}: {e}")
    except Exception as e:
        print(f"Erro ao abrir PDF com pdfplumber: {e}")

    # Fallback: PyMuPDF para PDFs com imagens embutidas
    try:
        doc = fitz.open(file_path)
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            image = Image.open(BytesIO(pix.tobytes()))
            enhanced_image = enhance_image(image)

            try:
                ocr_text = pytesseract.image_to_string(enhanced_image, lang='por')
                all_text += f"\nImagem na página {page_num + 1}:\n{ocr_text}\n"
            except Exception as ocr_error:
                print(f"Erro ao realizar OCR na página {page_num + 1}: {ocr_error}")
    except Exception as e:
        print(f"Erro ao processar imagens do PDF: {e}")

    return all_text, True, None

